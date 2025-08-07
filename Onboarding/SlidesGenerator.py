import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pygments import highlight
from pygments.lexers import PythonLexer
from pygments.formatters import HtmlFormatter
from bs4 import BeautifulSoup

# Load JSON
with open(input("JSON File Name (Must be in same level of directory) > "), "r", encoding="utf-8") as f:
    slides_data = json.load(f)

# Create Presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)  # 16:9 ratio

# === CONFIG ===
CODE_FONT = 'Courier New'
CODE_FONT_SIZE = Pt(12)
TEXT_FONT_SIZE = Pt(18)
LIGHT_GRAY = RGBColor(240, 240, 240)
DARK_GRAY = RGBColor(40, 40, 40)
TEXT_COLOR = RGBColor(0, 0, 0)

# === HELPERS ===
def hex_to_rgb(hex_color):
    hex_color = hex_color.strip('#')
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)

def extract_code_blocks(md):
    return re.findall(r"```[a-z]*\n(.*?)```", md, re.DOTALL)

def remove_code_blocks(md):
    return re.sub(r"```[a-z]*\n.*?```", '', md, flags=re.DOTALL).strip()

def is_markdown_table(md):
    return bool(re.search(r'\|.+\|\n\|[-: ]+\|\n\|?.+\|?', md.strip(), re.DOTALL))

def markdown_to_bullets(text):
    return [line.strip()[1:].strip() for line in text.strip().splitlines() if line.strip().startswith(("*", "-"))]

def add_explain_box(slide, text, left, top, width, height):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    if text.strip().startswith("*"):
        for bullet in markdown_to_bullets(text):
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = TEXT_FONT_SIZE
            p.level = 0
    else:
        tf.text = text
        tf.paragraphs[0].font.size = TEXT_FONT_SIZE

def add_code_box(slide, code, left, top, width, height):
    lexer = PythonLexer()
    formatter = HtmlFormatter(nowrap=True, style='default')
    highlighted = highlight(code, lexer, formatter)
    soup = BeautifulSoup(highlighted, 'html.parser')

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for span in soup.find_all(['span']):
        text = span.get_text()
        if not text.strip():
            continue
        color = span.get('style')
        color_match = re.search(r'color:\s*#([0-9a-fA-F]{6})', color or '')
        rgb = hex_to_rgb(color_match.group(1)) if color_match else RGBColor(0, 0, 0)

        p = tf.add_paragraph()
        p.text = text
        p.font.name = CODE_FONT
        p.font.size = CODE_FONT_SIZE
        p.font.color.rgb = rgb

    box.fill.solid()
    box.fill.fore_color.rgb = DARK_GRAY

def add_markdown_table(slide, markdown, left, top, width, height):
    lines = markdown.strip().splitlines()
    headers = [cell.strip() for cell in lines[0].split('|') if cell.strip()]
    rows = [line for line in lines[2:] if '|' in line]

    table_shape = slide.shapes.add_table(len(rows)+1, len(headers), left, top, width, height)
    table = table_shape.table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = TEXT_FONT_SIZE
        cell.text_frame.paragraphs[0].font.bold = True

    for r, row_text in enumerate(rows):
        cells = [cell.strip() for cell in row_text.split('|') if cell.strip()]
        for c, value in enumerate(cells):
            cell = table.cell(r+1, c)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = TEXT_FONT_SIZE

# === Slide Creator ===
def create_slide(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title at top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.text = data.get("title", "")
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True

    # Layout: Left (Explain), Right (Content)
    explain_area = (Inches(0.5), Inches(1), Inches(4), Inches(4))
    content_area = (Inches(5.1), Inches(1), Inches(4.4), Inches(4))

    if "explain" in data:
        add_explain_box(slide, data["explain"], *explain_area)

    if "content" in data:
        content = data["content"]

        if is_markdown_table(content):
            add_markdown_table(slide, content, *content_area)
        else:
            code_blocks = extract_code_blocks(content)
            if code_blocks:
                add_code_box(slide, "\n".join(code_blocks), *content_area)
            else:
                add_explain_box(slide, content, *content_area)

# === Build Slides ===
for slide_data in slides_data:
    create_slide(slide_data)

file_name = input("PPTX File Name > ")

prs.save(file_name)
print(f"Saved as {file_name}")
